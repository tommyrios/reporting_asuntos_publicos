from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from config import (
    BBVA_BLUE,
    BBVA_DARK_TEXT,
    BBVA_GREY_4,
    BBVA_GREY_2,
    BBVA_LIME,
    BBVA_MIDNIGHT,
    BBVA_SAND,
    BBVA_SERENE_BLUE,
    BBVA_WHITE,
    BODY_FONT,
    BRAND_DIR,
    TITLE_FONT,
)

LOGO_BLUE = BRAND_DIR / "bbva_logo_blue.png"
LOGO_WHITE = BRAND_DIR / "bbva_logo_white.png"
EMOJI_POLITICO = BRAND_DIR / "emoji_politico.png"

SLIDE_W = 13.333
SLIDE_H = 7.5

TEXT_LIMITS = {
    "headline": 86,
    "executive_vision": 430,
    "news_title": 82,
    "news_summary": 190,
    "news_meta": 68,
    "watchlist": 105,
}


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def _set_background(slide, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


def _set_font(
    run,
    font_name: str,
    size: float,
    color: str,
    bold: bool = False,
    italic: bool = False,
) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)


def _add_logo(slide, color: str, x: float, y: float, width: float) -> None:
    logo = LOGO_WHITE if color == "white" else LOGO_BLUE
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(x), Inches(y), width=Inches(width))
        return

    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "BBVA"
    _set_font(run, BODY_FONT, 20, BBVA_WHITE if color == "white" else BBVA_BLUE, bold=True)


def _add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: float = 18,
    color: str = BBVA_DARK_TEXT,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
    font_name: str = BODY_FONT,
    auto_size: bool = False,
    margin: float = 0.02,
) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = vertical_anchor
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    _set_font(run, font_name, font_size, color, bold=bold, italic=italic)


def _add_round_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    line: str | None = None,
) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line or fill)
    shape.line.width = Pt(0.5)
    return shape


def _add_separator(slide, x: float, y: float, w: float, color: str, height: float = 0.012) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(height))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(color)
    line.line.color.rgb = _rgb(color)


def _normalize_text(value: Any) -> str:
    return re.sub(r"\s+", " ", str(value or "")).strip()


def _sentence_safe_clip(value: Any, max_chars: int) -> str:
    """Clip at sentence boundary. Never returns visible ellipsis."""
    text = _normalize_text(value)
    if len(text) <= max_chars:
        return text

    sentences = re.split(r"(?<=[.!?])\s+", text)
    kept: list[str] = []
    size = 0
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        candidate_size = size + len(sentence) + (1 if kept else 0)
        if candidate_size <= max_chars:
            kept.append(sentence)
            size = candidate_size
        else:
            break

    if kept:
        return " ".join(kept).strip()

    clipped = text[:max_chars].rsplit(" ", 1)[0].strip()
    if clipped and clipped[-1] not in ".!?:;":
        clipped += "."
    return clipped


def _clean_text(value: Any, limit: int | None = None, sentence_safe: bool = True) -> str:
    text = _normalize_text(value)
    if not limit or len(text) <= limit:
        return text
    if sentence_safe:
        return _sentence_safe_clip(text, limit)
    clipped = text[:limit].rsplit(" ", 1)[0].strip()
    if clipped and clipped[-1] not in ".!?:;":
        clipped += "."
    return clipped


def _period_label(report: dict[str, Any]) -> str:
    period = report.get("period", {}) if isinstance(report.get("period"), dict) else {}
    return str(period.get("label") or "Período quincenal")


def _format_date(value: Any) -> str:
    raw = str(value or "").strip()
    if not raw:
        return ""
    for candidate in (raw, raw.replace("Z", "+00:00")):
        try:
            return datetime.fromisoformat(candidate).strftime("%d/%m/%Y")
        except Exception:
            pass
    return raw[:10]


def _add_header(slide, title: str, section: str, period_label: str) -> None:
    _add_logo(slide, "blue", 11.78, 0.32, 0.82)
    _add_textbox(slide, "INFORME QUINCENAL / ASUNTOS PÚBLICOS", 0.55, 0.31, 3.4, 0.20, font_size=6.7, color=BBVA_BLUE, bold=True)
    _add_textbox(slide, section.upper(), 4.15, 0.31, 3.3, 0.20, font_size=6.7, color=BBVA_BLUE, bold=True)
    _add_textbox(slide, title, 0.55, 0.70, 7.8, 0.55, font_size=24.5, color=BBVA_BLUE, bold=True, font_name=TITLE_FONT, auto_size=True)
    _add_textbox(slide, period_label, 0.56, 1.18, 5.0, 0.22, font_size=8.5, color=BBVA_GREY_4)


def _add_page_number(slide, page: int) -> None:
    _add_textbox(slide, f"p. {page}", 12.35, 7.05, 0.45, 0.18, font_size=7.2, color=BBVA_BLUE, align=PP_ALIGN.RIGHT)


def _add_cover_micro_fallback(slide) -> None:
    _add_round_rect(slide, 10.75, 5.18, 0.40, 1.28, fill=BBVA_SERENE_BLUE, line=BBVA_SERENE_BLUE)
    _add_round_rect(slide, 11.28, 4.82, 0.40, 1.64, fill=BBVA_SERENE_BLUE, line=BBVA_SERENE_BLUE)
    _add_round_rect(slide, 11.82, 5.45, 0.40, 1.01, fill=BBVA_SERENE_BLUE, line=BBVA_SERENE_BLUE)


def _add_cover_icon(slide) -> None:
    if EMOJI_POLITICO.exists():
        # Bottom-right micro asset. Small enough to work as a brand resource, not as the visual protagonist.
        slide.shapes.add_picture(str(EMOJI_POLITICO), Inches(10.50), Inches(4.72), width=Inches(2.05))
    else:
        _add_cover_micro_fallback(slide)


def _add_cover(slide, period_label: str) -> None:
    _set_background(slide, BBVA_BLUE)
    _add_logo(slide, "white", 0.55, 0.42, 1.02)

    _add_textbox(slide, period_label, 0.58, 2.78, 4.8, 0.22, font_size=9.8, color=BBVA_WHITE)
    _add_textbox(slide, "Dirección de Relaciones Institucionales", 0.58, 3.04, 5.9, 0.25, font_size=9.8, color=BBVA_WHITE)
    _add_separator(slide, 0.58, 3.42, 12.15, BBVA_WHITE, height=0.010)

    _add_textbox(
        slide,
        "Contexto político",
        0.55,
        3.90,
        7.2,
        1.05,
        font_size=30,
        color=BBVA_WHITE,
        bold=True,
        font_name=TITLE_FONT,
        auto_size=True,
        margin=0.00,
    )
    _add_textbox(slide, "Informe quincenal · borrador editable", 0.58, 6.92, 5.6, 0.22, font_size=7.8, color=BBVA_WHITE, italic=True)
    _add_cover_icon(slide)


def _paragraph(tf, text: str, font_size: float, color: str, bold: bool = False, font_name: str = BODY_FONT):
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    _set_font(run, font_name, font_size, color, bold=bold)
    return p


def _add_vision_card(slide, title: str, body: str, x: float, y: float, w: float, h: float) -> None:
    _add_round_rect(slide, x, y, w, h, fill=BBVA_WHITE, line=BBVA_WHITE)
    _add_textbox(slide, title, x + 0.24, y + 0.15, w - 0.48, 0.25, font_size=10.6, color=BBVA_BLUE, bold=True)
    _add_separator(slide, x + 0.24, y + 0.50, w - 0.48, BBVA_SERENE_BLUE, height=0.014)
    shape = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.66), Inches(w - 0.50), Inches(h - 0.80))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    text = _clean_text(body, TEXT_LIMITS["executive_vision"], sentence_safe=True)
    _paragraph(tf, text, 8.3, BBVA_DARK_TEXT, font_name=BODY_FONT)


def _add_news_card(slide, item: dict[str, Any], idx: int, x: float, y: float, w: float, h: float) -> None:
    _add_round_rect(slide, x, y, w, h, fill=BBVA_WHITE, line=BBVA_WHITE)

    number = f"{idx:02d}"
    _add_round_rect(slide, x + 0.20, y + 0.20, 0.42, 0.30, fill=BBVA_BLUE, line=BBVA_BLUE)
    _add_textbox(slide, number, x + 0.20, y + 0.245, 0.42, 0.12, font_size=6.5, color=BBVA_WHITE, bold=True, align=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)

    title = _clean_text(item.get("title", ""), TEXT_LIMITS["news_title"], sentence_safe=True)
    _add_textbox(slide, title, x + 0.73, y + 0.18, w - 0.95, 0.40, font_size=8.4, color=BBVA_BLUE, bold=True, auto_size=True)
    _add_separator(slide, x + 0.22, y + 0.68, w - 0.44, BBVA_GREY_2, height=0.010)

    source = _normalize_text(item.get("source", ""))
    date = _format_date(item.get("date") or item.get("published_at"))
    meta = " · ".join([part for part in [source, date] if part])
    if meta:
        _add_textbox(slide, _clean_text(meta, TEXT_LIMITS["news_meta"], sentence_safe=False), x + 0.24, y + 0.78, w - 0.48, 0.18, font_size=6.4, color=BBVA_GREY_4, bold=True)

    summary = item.get("why_it_matters") or item.get("summary") or item.get("reading") or ""
    summary = _clean_text(summary, TEXT_LIMITS["news_summary"], sentence_safe=True)
    shape = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 1.02), Inches(w - 0.48), Inches(h - 1.16))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    _paragraph(tf, summary, 7.2, BBVA_DARK_TEXT, font_name=BODY_FONT)


def _top_news_from_report(report: dict[str, Any]) -> list[dict[str, Any]]:
    analysis = report.get("analysis", {}) if isinstance(report.get("analysis"), dict) else {}
    top_news = analysis.get("top_political_news")
    if isinstance(top_news, list) and top_news:
        cleaned: list[dict[str, Any]] = []
        for item in top_news[:4]:
            if isinstance(item, dict):
                cleaned.append(item)
            elif item:
                cleaned.append({"title": str(item), "why_it_matters": ""})
        if cleaned:
            return cleaned

    selected_news = report.get("selected_news", []) if isinstance(report.get("selected_news"), list) else []
    fallback: list[dict[str, Any]] = []
    for item in selected_news[:4]:
        fallback.append(
            {
                "title": item.get("title", "Noticia política relevante"),
                "source": item.get("source", ""),
                "published_at": item.get("published_at", ""),
                "why_it_matters": item.get("summary", ""),
            }
        )
    return fallback


def _add_analysis(slide, report: dict[str, Any]) -> None:
    analysis = report.get("analysis", {}) if isinstance(report.get("analysis"), dict) else {}
    period_label = _period_label(report)
    stats = report.get("stats", {}) if isinstance(report.get("stats"), dict) else {}
    generated_at = datetime.now().strftime("%d/%m/%Y")

    _set_background(slide, BBVA_SAND)
    _add_header(slide, "Actualidad política", "Contexto político", period_label)

    headline = _clean_text(analysis.get("headline") or "Cuatro claves políticas de la quincena", TEXT_LIMITS["headline"], sentence_safe=True)
    _add_textbox(slide, headline, 0.56, 1.45, 11.5, 0.34, font_size=13.2, color=BBVA_DARK_TEXT, bold=True, auto_size=True)

    _add_vision_card(
        slide,
        "Visión ejecutiva",
        analysis.get("executive_vision", ""),
        0.55,
        1.88,
        12.20,
        1.32,
    )

    top_news = _top_news_from_report(report)
    card_positions = [
        (0.55, 3.46, 5.92, 1.55),
        (6.82, 3.46, 5.93, 1.55),
        (0.55, 5.22, 5.92, 1.55),
        (6.82, 5.22, 5.93, 1.55),
    ]
    for idx, (x, y, w, h) in enumerate(card_positions, start=1):
        item = top_news[idx - 1] if idx - 1 < len(top_news) else {"title": "Noticia pendiente de validación", "why_it_matters": "No se identificó una cuarta noticia con suficiente cobertura para esta quincena."}
        _add_news_card(slide, item, idx, x, y, w, h)

    footer = (
        f"Fuentes relevadas: {stats.get('raw_news_count', 0)} | "
        f"Noticias seleccionadas: {stats.get('selected_news_count', 0)} | "
        f"Generado: {generated_at}"
    )
    _add_textbox(slide, footer, 0.55, 6.95, 9.7, 0.18, font_size=7.0, color=BBVA_GREY_4)
    _add_page_number(slide, 2)


def _add_closing(slide) -> None:
    _set_background(slide, BBVA_BLUE)
    # Smaller logo to reduce pixelation and match the clean closing proportion.
    _add_logo(slide, "white", 5.98, 3.29, 1.36)


def create_aapp_pptx(report: dict[str, Any], output_path: Path) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    period_label = _period_label(report)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    _add_cover(slide, period_label)

    slide = prs.slides.add_slide(blank)
    _add_analysis(slide, report)

    slide = prs.slides.add_slide(blank)
    _add_closing(slide)

    prs.save(str(output_path))
    return output_path
